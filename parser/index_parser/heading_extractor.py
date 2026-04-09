"""원본 파일에서 Heading 스타일 기반 TOC 직접 추출 (LLM 불필요)

지원:
  - .docx : python-pptx Heading 1/2/3 스타일
  - .pptx : 슬라이드 제목(title placeholder) → level 1
"""
from __future__ import annotations
import re
from pathlib import Path


def extract_from_docx(docx_path: str) -> list[dict]:
    """docx Heading 스타일에서 계층 구조 TOC 추출.

    Returns:
        [{"title": ..., "level": 1|2|3, "page": None, "children": []}, ...]
        Heading이 하나도 없으면 빈 리스트 반환.
    """
    from pptx import Presentation  # noqa – pptx 미사용, docx용 별도 import
    from docx import Document

    doc = Document(docx_path)
    flat: list[dict] = []

    for para in doc.paragraphs:
        style_name = para.style.name if para.style else ""
        level = _heading_level(style_name)
        if level is None:
            continue
        title = para.text.strip()
        if not title or _is_toc_heading(title):
            continue
        flat.append({"title": title, "level": level, "page": None, "children": []})

    return _build_tree(flat) if flat else []


def extract_from_pptx(pptx_path: str) -> list[dict]:
    """PPTX 슬라이드 제목을 level 1 TOC 항목으로 추출."""
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN  # noqa

    prs = Presentation(pptx_path)
    flat: list[dict] = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.shape_type == 13:
                continue
            if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                ph_idx = shape.placeholder_format.idx
                if ph_idx == 0 and shape.has_text_frame:  # 제목 placeholder
                    title = shape.text_frame.text.strip()
                    if title:
                        flat.append({"title": title, "level": 1, "page": slide_idx, "children": []})
                    break

    return flat


# ── 내부 유틸 ─────────────────────────────────────────────

_TOC_HEADING_PATTERN = re.compile(r"^(목\s*차|table\s+of\s+contents|contents)$", re.IGNORECASE)

def _is_toc_heading(title: str) -> bool:
    return bool(_TOC_HEADING_PATTERN.match(title.strip()))


def _heading_level(style_name: str) -> int | None:
    """스타일명에서 Heading 레벨 추출. 비해당이면 None."""
    m = re.match(r"[Hh]eading\s+(\d)", style_name)
    if m:
        level = int(m.group(1))
        return min(level, 3)
    # 한글 워드 스타일: '제목 1', '제목 2'
    m2 = re.match(r"제목\s*(\d)", style_name)
    if m2:
        level = int(m2.group(1))
        return min(level, 3)
    return None


def _build_tree(flat: list[dict]) -> list[dict]:
    """평탄한 리스트를 children 계층 구조로 변환."""
    root: list[dict] = []
    stack: list[dict] = []  # (level, node)

    for item in flat:
        node = {**item, "children": []}
        level = node["level"]

        # 현재 level보다 크거나 같은 항목을 스택에서 제거
        while stack and stack[-1]["level"] >= level:
            stack.pop()

        if stack:
            stack[-1]["node"]["children"].append(node)
        else:
            root.append(node)

        stack.append({"level": level, "node": node})

    return root
