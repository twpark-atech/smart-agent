"""PPTX 파싱 - 텍스트/이미지/표 추출 (bbox 포함)"""
import io
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from .models import Block
from .llm import describe_image, describe_table
from .minio_helper import upload_image

# EMU → 포인트 변환 (1pt = 12700 EMU)
_EMU_TO_PT = 1 / 12700


def _emu_bbox(shape) -> list[float]:
    return [
        shape.left * _EMU_TO_PT,
        shape.top * _EMU_TO_PT,
        (shape.left + shape.width) * _EMU_TO_PT,
        (shape.top + shape.height) * _EMU_TO_PT,
    ]


def parse(pptx_path: str) -> list[Block]:
    """PPTX에서 슬라이드별 텍스트/이미지/표 추출."""
    blocks: list[Block] = []
    prs = Presentation(pptx_path)

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            bbox = _emu_bbox(shape)

            # 텍스트
            if shape.has_text_frame:
                text = "\n".join(
                    para.text for para in shape.text_frame.paragraphs if para.text.strip()
                )
                if text.strip():
                    blocks.append(Block(
                        block_type="text",
                        content=text,
                        page=slide_idx,
                        bbox=bbox,
                    ))

            # 표
            elif shape.has_table:
                table = shape.table
                headers = [table.cell(0, c).text for c in range(table.columns.__len__())]
                rows = []
                for r in range(1, table.rows.__len__()):
                    rows.append({
                        headers[c]: table.cell(r, c).text
                        for c in range(table.columns.__len__())
                    })
                md_lines = ["| " + " | ".join(headers) + " |"]
                md_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                for row in rows:
                    md_lines.append("| " + " | ".join(str(row.get(h, "")) for h in headers) + " |")
                table_md = "\n".join(md_lines)

                description = describe_table(table_md)
                blocks.append(Block(
                    block_type="table",
                    content=f"[표 설명]\n{description}\n\n[표 원문]\n{table_md}",
                    page=slide_idx,
                    bbox=bbox,
                    table_json=rows,
                ))

            # 이미지
            elif shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image = shape.image
                ext = image.ext  # "png", "jpeg" 등
                with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
                    tmp.write(image.blob)
                    tmp_path = tmp.name

                try:
                    minio_key = upload_image(tmp_path, source_pdf=pptx_path, page=slide_idx)
                    description = describe_image(tmp_path)
                finally:
                    Path(tmp_path).unlink(missing_ok=True)

                blocks.append(Block(
                    block_type="image",
                    content=f"[이미지 설명]\n{description}",
                    page=slide_idx,
                    bbox=bbox,
                    minio_key=minio_key,
                ))

    return blocks
